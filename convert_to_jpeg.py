#!/usr/bin/env python3
"""
Image to JPEG Converter — Full-featured edition
Features: Drag & drop, threading, thumbnails, reorder, dark mode, and more.
"""

import os
import sys
import shutil
import tempfile
import zipfile
import threading
import subprocess
import platform
from pathlib import Path

import tkinter as tk
from tkinter import filedialog, messagebox, ttk


# ── Auto-install dependencies ────────────────────────────────────────────────
def _ensure(*packages):
    import importlib
    for pkg, imp in packages:
        try:
            importlib.import_module(imp)
        except ImportError:
            print(f"Installing {pkg}…")
            subprocess.check_call([sys.executable, "-m", "pip", "install", pkg])

_ensure(
    ("Pillow", "PIL"),
    ("pymupdf", "fitz"),
    ("python-docx", "docx"),
    ("python-pptx", "pptx"),
    ("openpyxl", "openpyxl"),
    ("tkinterdnd2", "tkinterdnd2"),
)

from PIL import Image, ImageTk, ImageDraw
import fitz
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from tkinterdnd2 import TkinterDnD, DND_FILES


# ── Formats & constants ──────────────────────────────────────────────────────
IMAGE_FORMATS = (
    ("Image Files", "*.png *.jpg *.jpeg *.bmp *.gif *.tiff *.tif *.webp *.ico *.ppm *.tga *.heic *.heif"),
    ("All Files", "*.*"),
)

IMPORT_FORMATS = (
    ("All Supported Files", "*.pdf *.docx *.pptx *.xlsx *.zip *.cbz"),
    ("PDF", "*.pdf"),
    ("Word", "*.docx"),
    ("PowerPoint", "*.pptx"),
    ("Excel", "*.xlsx"),
    ("Archive", "*.zip *.cbz"),
    ("All Files", "*.*"),
)

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff", ".tif",
                    ".webp", ".ico", ".ppm", ".tga", ".heic", ".heif"}

DOC_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".zip", ".cbz"}

THEMES = {
    "light": {
        "bg":         "#f7fafc", "fg":         "#2d3748",
        "panel":      "#ffffff", "panel_fg":   "#2d3748",
        "entry_bg":   "#ffffff", "entry_fg":   "#2d3748",
        "list_bg":    "#ffffff", "list_fg":    "#2d3748", "list_sel": "#ebf4f0",
        "btn_bg":     "#edf2f7", "btn_fg":     "#4a5568",
        "hint":       "#718096",
        "accent":     "#2d3748", "accent_fg":  "#ffffff",
        "accent2":    "#48bb78", "accent3":    "#38a169",
        "danger":     "#e53e3e",
        "header_bg":  "#2d3748", "header_fg":  "#ffffff", "header_sub": "#a0aec0",
    },
    "dark": {
        "bg":         "#1a202c", "fg":         "#e2e8f0",
        "panel":      "#2d3748", "panel_fg":   "#e2e8f0",
        "entry_bg":   "#374151", "entry_fg":   "#e2e8f0",
        "list_bg":    "#2d3748", "list_fg":    "#e2e8f0", "list_sel": "#2f5a3d",
        "btn_bg":     "#4a5568", "btn_fg":     "#e2e8f0",
        "hint":       "#a0aec0",
        "accent":     "#4a5568", "accent_fg":  "#ffffff",
        "accent2":    "#48bb78", "accent3":    "#38a169",
        "danger":     "#fc4444",
        "header_bg":  "#171f2e", "header_fg":  "#e2e8f0", "header_sub": "#718096",
    },
}


# ── Core conversion ──────────────────────────────────────────────────────────
def convert_image_to_jpeg(input_path, output_dir, quality=90,
                          index=1, prefix="Image", used_names=None):
    try:
        input_path = Path(input_path)
        output_dir = Path(output_dir)
        output_path = output_dir / f"{prefix}_{index}.jpg"

        if used_names is not None:
            counter = index
            while str(output_path) in used_names or output_path.exists():
                counter += 1
                output_path = output_dir / f"{prefix}_{counter}.jpg"
            used_names.add(str(output_path))

        with Image.open(input_path) as img:
            if img.mode in ("RGBA", "LA", "P"):
                background = Image.new("RGB", img.size, (255, 255, 255))
                if img.mode == "P":
                    img = img.convert("RGBA")
                background.paste(img, mask=img.split()[-1] if img.mode in ("RGBA", "LA") else None)
                img = background
            elif img.mode != "RGB":
                img = img.convert("RGB")

            img.save(output_path, "JPEG", quality=quality, optimize=True)

        return True, f"✔  {input_path.name}  →  {output_path.name}"
    except Exception as e:
        return False, f"✘  {Path(input_path).name}  —  Error: {e}"


# ── Extraction helpers ───────────────────────────────────────────────────────
def extract_from_pdf(filepath, tmp_dir):
    extracted = []
    doc = fitz.open(filepath)
    for page_num in range(len(doc)):
        for img_index, img in enumerate(doc[page_num].get_images(full=True)):
            xref = img[0]
            base = doc.extract_image(xref)
            # Zero-padded so alphabetical sorting matches numeric order
            out = os.path.join(tmp_dir,
                               f"pdf_p{page_num+1:04d}_img{img_index+1:04d}.{base['ext']}")
            with open(out, "wb") as f:
                f.write(base["image"])
            extracted.append(out)
    doc.close()
    return extracted


def extract_from_docx(filepath, tmp_dir):
    """Extract images from DOCX in the order they appear in the document body,
    including duplicates (same image used multiple times)."""
    extracted = []
    doc = Document(filepath)

    BLIP_TAG = "{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
    EMBED_ATTR = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"

    idx = 1
    # Walk body XML in order — each <a:blip> is one image occurrence in the document
    for blip in doc.element.body.iter(BLIP_TAG):
        embed_id = blip.get(EMBED_ATTR)
        if not embed_id:
            continue
        if embed_id not in doc.part.rels:
            continue
        rel = doc.part.rels[embed_id]
        if "image" not in rel.reltype:
            continue
        ext = Path(rel.target_ref).suffix or ".png"
        # Each occurrence gets its own file — duplicates included
        out = os.path.join(tmp_dir, f"docx_img{idx:04d}{ext}")
        with open(out, "wb") as f:
            f.write(rel.target_part.blob)
        extracted.append(out)
        idx += 1
    return extracted


def extract_from_pptx(filepath, tmp_dir):
    extracted = []
    prs = Presentation(filepath)
    idx = 1
    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:
                ext = "." + shape.image.ext
                out = os.path.join(tmp_dir,
                                   f"pptx_s{slide_num+1:04d}_img{idx:04d}{ext}")
                with open(out, "wb") as f:
                    f.write(shape.image.blob)
                extracted.append(out)
                idx += 1
    return extracted


def extract_from_xlsx(filepath, tmp_dir):
    extracted = []
    wb = load_workbook(filepath)
    idx = 1
    for sheet in wb.worksheets:
        for img in sheet._images:
            try:
                out = os.path.join(tmp_dir, f"xlsx_img{idx:04d}.png")
                with open(out, "wb") as f:
                    f.write(img._data())
                extracted.append(out)
                idx += 1
            except Exception:
                pass
    return extracted


def extract_from_zip(filepath, tmp_dir):
    extracted = []
    with zipfile.ZipFile(filepath, "r") as zf:
        # Sort alphabetically so images come out in a predictable order
        for name in sorted(zf.namelist()):
            if Path(name).suffix.lower() in IMAGE_EXTENSIONS:
                out = os.path.join(tmp_dir, Path(name).name)
                with zf.open(name) as src, open(out, "wb") as dst:
                    dst.write(src.read())
                extracted.append(out)
    return extracted


def extract_images_from_file(filepath, tmp_dir):
    ext = Path(filepath).suffix.lower()
    try:
        if ext == ".pdf":
            imgs = extract_from_pdf(filepath, tmp_dir)
        elif ext == ".docx":
            imgs = extract_from_docx(filepath, tmp_dir)
        elif ext == ".pptx":
            imgs = extract_from_pptx(filepath, tmp_dir)
        elif ext == ".xlsx":
            imgs = extract_from_xlsx(filepath, tmp_dir)
        elif ext in (".zip", ".cbz"):
            imgs = extract_from_zip(filepath, tmp_dir)
        else:
            return [], f"✘  Unsupported file type: {ext}"

        if not imgs:
            return [], f"⚠  No images found in {Path(filepath).name}"
        return imgs, f"📦  {Path(filepath).name}  →  {len(imgs)} image(s) extracted"
    except Exception as e:
        return [], f"✘  {Path(filepath).name}  —  Error: {e}"


def open_folder(path):
    """Cross-platform open folder in file manager."""
    try:
        system = platform.system()
        if system == "Windows":
            os.startfile(path)
        elif system == "Darwin":
            subprocess.run(["open", path])
        else:
            subprocess.run(["xdg-open", path])
    except Exception as e:
        messagebox.showerror("Cannot open folder", str(e))


def shade(hex_color, factor):
    """Lighten (factor > 1) or darken (factor < 1) a hex color."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = max(0, min(255, int(r * factor)))
    g = max(0, min(255, int(g * factor)))
    b = max(0, min(255, int(b * factor)))
    return f"#{r:02x}{g:02x}{b:02x}"


# ── GUI ──────────────────────────────────────────────────────────────────────
class ConverterApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Image → JPEG Converter")
        self.root.geometry("620x680")
        self.root.minsize(560, 560)

        self._tmp_dirs = []
        self._tree_to_path = {}        # treeview item id → filepath
        self._thumbnails = []          # keep refs to prevent GC
        self._cancel_event = threading.Event()
        self._converting = False
        self._last_output_folder = None
        self._theme = "light"

        self._build_ui()
        self._apply_theme()
        self.root.protocol("WM_DELETE_WINDOW", self._on_close)

    # ── UI construction ──────────────────────────────────────────────────────

    def _build_ui(self):
        PAD = 10

        # Pre-generate custom button icons (perfect alignment vs. emoji rendering)
        self._folder_icon_white = self._make_folder_icon("white")

        # ── Header (truly centered title; theme button via place() on the right) ─
        self.header = tk.Frame(self.root, pady=14)
        self.header.pack(side="top", fill="x")

        self.title_label = tk.Label(self.header, text="Image  →  JPEG Converter",
                                    font=("Helvetica", 16, "bold"))
        self.title_label.pack(pady=(2, 4))

        self.subtitle = tk.Label(self.header,
                                 text="Add images or import files — convert anything to JPEG",
                                 font=("Helvetica", 9))
        self.subtitle.pack(pady=(0, 2))

        self.theme_btn = tk.Button(self.header, text="🌙  Dark", command=self._toggle_theme,
                                   relief="flat", padx=14, pady=5, cursor="hand2",
                                   font=("Helvetica", 9, "bold"), borderwidth=0)
        # place() positions absolutely without affecting pack layout — title stays centered
        self.theme_btn.place(relx=1.0, rely=0.5, anchor="e", x=-15)

        self.info_btn = tk.Button(self.header, text="ℹ", command=self._show_info,
                                  relief="flat", padx=10, pady=5, cursor="hand2",
                                  font=("Helvetica", 12, "bold"), borderwidth=0)
        self.info_btn.place(relx=0.0, rely=0.5, anchor="w", x=15)

        # ── Action buttons (anchored to bottom — packed FIRST so they stay) ──
        action_row = tk.Frame(self.root)
        action_row.pack(side="bottom", fill="x", padx=PAD, pady=(8, PAD))

        self.convert_btn = tk.Button(action_row, text="Convert to JPEG",
                                     command=self._convert,
                                     font=("Helvetica", 12, "bold"),
                                     relief="flat", pady=10, cursor="hand2", borderwidth=0)
        self.convert_btn.pack(side="left", padx=(0, 6), fill="x", expand=True)

        self.cancel_btn = tk.Button(action_row, text="✕  Cancel", command=self._cancel_convert,
                                    font=("Helvetica", 10, "bold"),
                                    relief="flat", pady=10, padx=14, cursor="hand2",
                                    borderwidth=0, state="disabled",
                                    disabledforeground="#ffe6e2")
        self.cancel_btn.pack(side="left", padx=(0, 6))

        self.open_folder_btn = tk.Button(action_row, text="  Open Output",
                                         command=self._open_output_folder,
                                         image=self._folder_icon_white, compound="left",
                                         font=("Helvetica", 10, "bold"),
                                         relief="flat", pady=10, padx=14, cursor="hand2",
                                         borderwidth=0, state="disabled",
                                         disabledforeground="#d6f5e1")
        self.open_folder_btn.pack(side="left")

        # ── Progress & log (above action row) ────────────────────────────────
        self.log_frame = tk.LabelFrame(self.root, text=" Progress ", padx=PAD, pady=6)
        self.log_frame.pack(side="bottom", fill="x", padx=PAD, pady=4)

        self.progress = ttk.Progressbar(self.log_frame, mode="determinate")
        self.progress.pack(fill="x", pady=(0, 4))

        log_container = tk.Frame(self.log_frame)
        log_container.pack(fill="both", expand=True)
        log_scroll = ttk.Scrollbar(log_container)
        log_scroll.pack(side="right", fill="y")
        self.log_box = tk.Listbox(log_container, height=4,
                                  yscrollcommand=log_scroll.set, font=("Courier", 9))
        self.log_box.pack(side="left", fill="both", expand=True)
        log_scroll.config(command=self.log_box.yview)

        # ── Quality (above log) ───────────────────────────────────────────────
        self.q_frame = tk.LabelFrame(self.root, text=" Image Settings ", padx=PAD, pady=4)
        self.q_frame.pack(side="bottom", fill="x", padx=PAD, pady=2)

        q_row = tk.Frame(self.q_frame)
        q_row.pack(fill="x")
        tk.Label(q_row, text="Quality:", width=10, anchor="w").pack(side="left")
        self.quality_var = tk.IntVar(value=90)
        self.q_scale = tk.Scale(q_row, from_=10, to=100, orient="horizontal", variable=self.quality_var,
                                length=300, showvalue=False, highlightthickness=0)
        self.q_scale.pack(side="left", fill="x", expand=True)
        self.q_label = tk.Label(q_row, text="90 (High)", width=14, anchor="w")
        self.q_label.pack(side="left", padx=6)
        self.quality_var.trace_add("write", self._update_q_label)

        # ── Output folder + Prefix (above quality) ───────────────────────────
        self.out_frame = tk.LabelFrame(self.root, text=" Output Settings ", padx=PAD, pady=4)
        self.out_frame.pack(side="bottom", fill="x", padx=PAD, pady=2)

        out_row = tk.Frame(self.out_frame)
        out_row.pack(fill="x")
        tk.Label(out_row, text="Folder:", width=10, anchor="w").pack(side="left")
        self.out_var = tk.StringVar(value="")
        self.out_entry = tk.Entry(out_row, textvariable=self.out_var, font=("Helvetica", 9))
        self.out_entry.insert(0, "No folder selected — click Browse to choose…")
        self.out_entry.config(fg="#999")
        self.out_entry.pack(side="left", fill="x", expand=True, padx=(0, 6))
        self.out_browse_btn = tk.Button(out_row, text="Browse…", command=self._browse_output,
                                        relief="flat", padx=8, cursor="hand2")
        self.out_browse_btn.pack(side="left", padx=(0, 4))
        self.out_clear_btn = tk.Button(out_row, text="Clear", command=self._clear_output,
                                       relief="flat", padx=8, cursor="hand2")
        self.out_clear_btn.pack(side="left")

        prefix_row = tk.Frame(self.out_frame)
        prefix_row.pack(fill="x", pady=(4, 0))
        tk.Label(prefix_row, text="Prefix:", width=10, anchor="w").pack(side="left")
        self.prefix_var = tk.StringVar(value="Image")
        self.prefix_entry = tk.Entry(prefix_row, textvariable=self.prefix_var,
                                     font=("Helvetica", 9), width=18)
        self.prefix_entry.pack(side="left", padx=(0, 6))
        self.prefix_hint = tk.Label(prefix_row, text="→  Image_1.jpg, Image_2.jpg, …",
                                    font=("Helvetica", 8))
        self.prefix_hint.pack(side="left")
        self.prefix_var.trace_add("write", self._update_prefix_hint)

        # ── File list with thumbnails (takes remaining space) ────────────────
        self.file_frame = tk.LabelFrame(self.root, text=" Selected Images ", padx=PAD, pady=8,
                                        font=("Helvetica", 9, "bold"))
        self.file_frame.pack(fill="both", expand=True, padx=PAD, pady=(0, 2))

        # IMPORTANT: pack button rows FIRST with side="bottom" so they always stay visible.
        # The tree container packs last with expand=True and shrinks to fit remaining space.

        self.hint_label = tk.Label(self.file_frame,
                                   text="💡 Drag & drop files/folders directly onto the list",
                                   font=("Helvetica", 8, "italic"))
        self.hint_label.pack(side="bottom", anchor="w", pady=(3, 0))

        # Single button row: primary actions on left, management buttons + count on right
        btn_row = tk.Frame(self.file_frame)
        btn_row.pack(side="bottom", fill="x", pady=(6, 2))

        # Left side: primary big buttons
        self.add_btn = tk.Button(btn_row, text="➕  Add Images", command=self._add_files,
                                 relief="flat", padx=12, pady=7, cursor="hand2",
                                 font=("Helvetica", 10, "bold"), borderwidth=0)
        self.add_btn.pack(side="left", padx=(0, 5))

        self.import_btn = tk.Button(btn_row, text="  Import File", command=self._import_file,
                                    image=self._folder_icon_white, compound="left",
                                    relief="flat", padx=12, pady=7, cursor="hand2",
                                    font=("Helvetica", 10, "bold"), borderwidth=0)
        self.import_btn.pack(side="left", padx=(0, 5))

        # Right side: count label + management buttons (packed right-to-left visually)
        self.count_label = tk.Label(btn_row, text="0 file(s)", font=("Helvetica", 9, "bold"))
        self.count_label.pack(side="right", padx=(6, 0))

        self.clear_btn = tk.Button(btn_row, text="🗑 Clear", command=self._clear_files,
                                   relief="flat", padx=6, pady=4, cursor="hand2",
                                   font=("Helvetica", 9), borderwidth=0)
        self.clear_btn.pack(side="right", padx=(2, 0))

        self.remove_btn = tk.Button(btn_row, text="✕ Remove", command=self._remove_selected,
                                    relief="flat", padx=6, pady=4, cursor="hand2",
                                    font=("Helvetica", 9), borderwidth=0)
        self.remove_btn.pack(side="right", padx=(2, 0))

        self.down_btn = tk.Button(btn_row, text="↓ Down", command=self._move_down,
                                  relief="flat", padx=6, pady=4, cursor="hand2",
                                  font=("Helvetica", 9), borderwidth=0)
        self.down_btn.pack(side="right", padx=(2, 0))

        self.up_btn = tk.Button(btn_row, text="↑ Up", command=self._move_up,
                                relief="flat", padx=6, pady=4, cursor="hand2",
                                font=("Helvetica", 9), borderwidth=0)
        self.up_btn.pack(side="right", padx=(2, 0))

        # Tree container — fills remaining space at the top
        tree_container = tk.Frame(self.file_frame)
        tree_container.pack(side="top", fill="both", expand=True)

        style = ttk.Style()
        style.configure("Thumbs.Treeview", rowheight=52)

        self.file_tree = ttk.Treeview(tree_container, columns=("path",), style="Thumbs.Treeview",
                                      selectmode="extended", height=3)
        self.file_tree.heading("#0", text="  Image")
        self.file_tree.heading("path", text="Location")
        self.file_tree.column("#0", width=300, minwidth=200, stretch=False)
        self.file_tree.column("path", width=550, minwidth=200, stretch=False)

        tree_yscroll = ttk.Scrollbar(tree_container, orient="vertical",
                                     command=self.file_tree.yview)
        tree_xscroll = ttk.Scrollbar(tree_container, orient="horizontal",
                                     command=self.file_tree.xview)
        self.file_tree.config(yscrollcommand=tree_yscroll.set, xscrollcommand=tree_xscroll.set)

        self.file_tree.grid(row=0, column=0, sticky="nsew")
        tree_yscroll.grid(row=0, column=1, sticky="ns")
        tree_xscroll.grid(row=1, column=0, sticky="ew")
        tree_container.rowconfigure(0, weight=1)
        tree_container.columnconfigure(0, weight=1)

        self.file_tree.drop_target_register(DND_FILES)
        self.file_tree.dnd_bind("<<Drop>>", self._on_drop)

    # ── Theme ────────────────────────────────────────────────────────────────

    def _toggle_theme(self):
        self._theme = "dark" if self._theme == "light" else "light"
        self.theme_btn.config(text="☀  Light" if self._theme == "dark" else "🌙  Dark")
        self._apply_theme()

    def _show_info(self):
        info = (
            "Image → JPEG Converter\n"
            "────────────────────────────────\n\n"
            "ADDING FILES\n"
            "  ➕ Add Images   — pick image files directly\n"
            "  📂 Import File  — extract images from:\n"
            "       PDF, DOCX, PPTX, XLSX, ZIP / CBZ\n"
            "  Drag & Drop     — files or entire folders\n\n"
            "FILE LIST\n"
            "  ↑ Up / ↓ Down  — reorder files\n"
            "                   (sets output numbering)\n"
            "  ✕ Remove       — remove selected files\n"
            "  🗑 Clear All   — remove all files\n"
            "  Ctrl+click     — select multiple files\n\n"
            "OUTPUT SETTINGS\n"
            "  Folder  — required, where JPEGs are saved\n"
            "  Prefix  — base name for output files\n"
            "            e.g. 'Photo' → Photo_1.jpg…\n\n"
            "IMAGE SETTINGS\n"
            "  Quality slider (10–100)\n"
            "    10–49  Low   |  50–74  Medium\n"
            "    75–94  High  |  95–100 Maximum\n"
            "  Default: 90 (High)\n\n"
            "CONVERSION\n"
            "  Convert to JPEG — starts the batch\n"
            "  Cancel          — stops after current file\n"
            "  Open Output     — opens output folder\n\n"
            "NOTES\n"
            "  • Transparent images (PNG, WEBP, GIF) are\n"
            "    composited onto a white background\n"
            "  • Duplicate images are kept — each gets\n"
            "    its own output file\n"
            "  • Temp files from Import are auto-deleted\n"
            "    when the app closes\n"
            "  • Window is fully resizable\n"
        )
        win = tk.Toplevel(self.root)
        win.title("App Info")
        win.resizable(False, False)
        t = THEMES[self._theme]
        win.config(bg=t["bg"])

        # Header strip
        hdr = tk.Frame(win, bg=t["header_bg"], pady=10)
        hdr.pack(fill="x")
        tk.Label(hdr, text="ℹ  Feature Guide",
                 font=("Helvetica", 13, "bold"),
                 bg=t["header_bg"], fg=t["header_fg"]).pack()

        # Scrollable text area
        txt_frame = tk.Frame(win, bg=t["bg"])
        txt_frame.pack(fill="both", expand=True, padx=14, pady=10)
        scrollbar = tk.Scrollbar(txt_frame)
        scrollbar.pack(side="right", fill="y")
        txt = tk.Text(txt_frame, font=("Courier", 9), wrap="word",
                      width=52, height=28,
                      bg=t["list_bg"], fg=t["list_fg"],
                      relief="flat", borderwidth=0,
                      yscrollcommand=scrollbar.set,
                      padx=10, pady=8, state="normal")
        txt.insert("1.0", info)
        txt.config(state="disabled")
        txt.pack(side="left", fill="both", expand=True)
        scrollbar.config(command=txt.yview)

        close_btn = tk.Button(win, text="Close", command=win.destroy,
                              font=("Helvetica", 10, "bold"),
                              bg=t["accent"], fg=t["accent_fg"],
                              relief="flat", pady=8, padx=20,
                              cursor="hand2", borderwidth=0)
        close_btn.pack(pady=(0, 12))
        win.transient(self.root)
        win.grab_set()

    def _apply_theme(self):
        t = THEMES[self._theme]

        self.root.config(bg=t["bg"])

        # ── Header: dedicated header colors (changes with theme) ─────────────
        self.header.config(bg=t["header_bg"])
        self.title_label.config(bg=t["header_bg"], fg=t["header_fg"])
        self.subtitle.config(bg=t["header_bg"], fg=t["header_sub"])
        self._style_btn(self.theme_btn,
                        bg=shade(t["header_bg"], 1.3 if self._theme == "light" else 1.5),
                        fg=t["header_fg"])
        self._style_btn(self.info_btn,
                        bg=shade(t["header_bg"], 1.3 if self._theme == "light" else 1.5),
                        fg=t["header_fg"])

        # ── Section frames ───────────────────────────────────────────────────
        for frame in (self.file_frame, self.out_frame, self.q_frame, self.log_frame):
            frame.config(bg=t["bg"], fg=t["fg"])

        for frame in (self.file_frame, self.out_frame, self.q_frame, self.log_frame):
            self._theme_children(frame, t)

        # ── Action buttons (with hover) ──────────────────────────────────────
        self._style_btn(self.convert_btn, bg=t["accent"], fg=t["accent_fg"])
        self._style_btn(self.cancel_btn, bg=t["danger"], fg="#ffffff",
                        disabled_fg="#ffe6e2")
        self._style_btn(self.open_folder_btn, bg=t["accent2"], fg="#ffffff",
                        disabled_fg="#d6f5e1")

        # ── File-panel action buttons ────────────────────────────────────────
        self._style_btn(self.add_btn, bg=t["accent"], fg=t["accent_fg"])
        self._style_btn(self.import_btn, bg=t["accent2"], fg="#ffffff")

        # ── Compact list buttons ─────────────────────────────────────────────
        for btn in (self.up_btn, self.down_btn, self.remove_btn, self.clear_btn,
                    self.out_browse_btn, self.out_clear_btn):
            self._style_btn(btn, bg=t["btn_bg"], fg=t["btn_fg"])

        # ── Log box ──────────────────────────────────────────────────────────
        self.log_box.config(bg=t["list_bg"], fg=t["list_fg"],
                            selectbackground=t["list_sel"], selectforeground=t["list_fg"])

        # ── Hints ────────────────────────────────────────────────────────────
        self.hint_label.config(fg=t["hint"])
        self.prefix_hint.config(fg=t["hint"])

        # ── Treeview style ───────────────────────────────────────────────────
        style = ttk.Style()
        style.theme_use("default")
        style.configure("Thumbs.Treeview",
                        background=t["list_bg"], foreground=t["list_fg"],
                        fieldbackground=t["list_bg"], rowheight=52,
                        borderwidth=0)
        style.map("Thumbs.Treeview",
                  background=[("selected", t["list_sel"])],
                  foreground=[("selected", t["list_fg"])])
        style.configure("Thumbs.Treeview.Heading",
                        background=t["btn_bg"], foreground=t["fg"],
                        font=("Helvetica", 9, "bold"),
                        relief="flat", borderwidth=0)

    def _style_btn(self, btn, bg, fg, disabled_fg=None):
        """Apply consistent button styling + hover transition."""
        hover_bg = shade(bg, 0.88 if self._theme == "light" else 1.18)
        active_bg = shade(bg, 0.78 if self._theme == "light" else 1.30)
        config = {
            "bg": bg, "fg": fg,
            "activebackground": active_bg, "activeforeground": fg,
        }
        if disabled_fg:
            config["disabledforeground"] = disabled_fg
        btn.config(**config)
        # Remove any previous hover bindings to avoid stacking
        btn.unbind("<Enter>")
        btn.unbind("<Leave>")
        btn.bind("<Enter>",
                 lambda e, b=btn, c=hover_bg: b.config(bg=c) if b["state"] != "disabled" else None)
        btn.bind("<Leave>",
                 lambda e, b=btn, c=bg: b.config(bg=c) if b["state"] != "disabled" else None)

    def _theme_children(self, widget, t):
        """Recursively apply theme to child tk widgets (skips buttons we style ourselves)."""
        # Buttons that get their own _style_btn treatment — don't override them here
        styled = {self.add_btn, self.import_btn, self.convert_btn,
                  self.cancel_btn, self.open_folder_btn, self.theme_btn,
                  self.info_btn, self.up_btn, self.down_btn, self.remove_btn,
                  self.clear_btn, self.out_browse_btn, self.out_clear_btn}
        for child in widget.winfo_children():
            cls = child.winfo_class()
            try:
                if cls in ("Frame", "Labelframe"):
                    child.config(bg=t["bg"])
                    if cls == "Labelframe":
                        child.config(fg=t["fg"])
                elif cls == "Label":
                    child.config(bg=t["bg"], fg=t["fg"])
                elif cls == "Button":
                    if child not in styled:
                        child.config(bg=t["btn_bg"], fg=t["btn_fg"],
                                     activebackground=t["btn_bg"], activeforeground=t["btn_fg"])
                elif cls == "Entry":
                    if child is self.out_entry and not self.out_var.get():
                        child.config(bg=t["entry_bg"], fg=t["hint"],
                                     insertbackground=t["fg"])
                    else:
                        child.config(bg=t["entry_bg"], fg=t["entry_fg"],
                                     insertbackground=t["fg"])
                elif cls == "Scale":
                    child.config(bg=t["bg"], fg=t["fg"], troughcolor=t["entry_bg"],
                                 highlightbackground=t["bg"], activebackground=t["accent"])
                elif cls == "Checkbutton":
                    child.config(bg=t["bg"], fg=t["fg"], selectcolor=t["entry_bg"],
                                 activebackground=t["bg"], activeforeground=t["fg"])
            except tk.TclError:
                pass
            self._theme_children(child, t)

    # ── Custom icon generation (for pixel-perfect button alignment) ──────────

    def _make_folder_icon(self, color="white", w=16, h=13):
        """Create a clean folder-shaped PhotoImage that aligns with text baseline."""
        img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        draw = ImageDraw.Draw(img)
        # Folder tab (small rect at top-left)
        draw.rectangle([(0, 2), (w // 2, 4)], fill=color)
        # Folder body (main rectangle)
        draw.rectangle([(0, 4), (w - 1, h - 2)], fill=color)
        return ImageTk.PhotoImage(img)

    # ── Thumbnail helpers ────────────────────────────────────────────────────

    def _make_thumb(self, filepath, size=44):
        try:
            img = Image.open(filepath)
            img.thumbnail((size, size), Image.LANCZOS)
            if img.mode not in ("RGB", "RGBA"):
                img = img.convert("RGBA")
            return ImageTk.PhotoImage(img)
        except Exception:
            return None

    # ── File management ──────────────────────────────────────────────────────

    def _add_files(self):
        files = filedialog.askopenfilenames(title="Select Images", filetypes=IMAGE_FORMATS)
        if files:
            self._enqueue_files(list(files))

    def _enqueue_files(self, files):
        for f in files:
            thumb = self._make_thumb(f)
            if thumb:
                self._thumbnails.append(thumb)
            item_id = self.file_tree.insert(
                "", "end",
                image=thumb if thumb else "",
                text="  " + Path(f).name,
                values=(str(Path(f).parent),),
            )
            self._tree_to_path[item_id] = f
        self._update_count()

    def _remove_selected(self):
        for item in self.file_tree.selection():
            if item in self._tree_to_path:
                del self._tree_to_path[item]
            self.file_tree.delete(item)
        self._update_count()

    def _clear_files(self):
        for item in self.file_tree.get_children():
            self.file_tree.delete(item)
        self._tree_to_path.clear()
        self._thumbnails.clear()
        self._update_count()

    def _move_up(self):
        for item in self.file_tree.selection():
            idx = self.file_tree.index(item)
            if idx > 0:
                self.file_tree.move(item, "", idx - 1)

    def _move_down(self):
        items = list(self.file_tree.selection())
        items.reverse()
        children = self.file_tree.get_children()
        for item in items:
            idx = self.file_tree.index(item)
            if idx < len(children) - 1:
                self.file_tree.move(item, "", idx + 1)

    def _get_ordered_files(self):
        return [self._tree_to_path[item] for item in self.file_tree.get_children()
                if item in self._tree_to_path]

    def _update_count(self):
        n = len(self.file_tree.get_children())
        self.count_label.config(text=f"{n} file(s)")

    # ── Drag & drop ──────────────────────────────────────────────────────────

    def _on_drop(self, event):
        # tkinterdnd2 returns a string with paths; splitlist handles spaces/braces
        raw = self.root.tk.splitlist(event.data)

        image_files = []
        doc_files = []

        def consider(path):
            ext = Path(path).suffix.lower()
            if ext in IMAGE_EXTENSIONS:
                image_files.append(path)
            elif ext in DOC_EXTENSIONS:
                doc_files.append(path)

        for raw_path in raw:
            p = raw_path.strip("{}")
            if os.path.isdir(p):
                for root_dir, _, names in os.walk(p):
                    for name in names:
                        consider(os.path.join(root_dir, name))
            elif os.path.isfile(p):
                consider(p)

        if image_files:
            self._enqueue_files(image_files)

        if doc_files:
            tmp_dir = tempfile.mkdtemp(prefix="jpeg_dnd_")
            self._tmp_dirs.append(tmp_dir)
            for f in doc_files:
                imgs, summary = extract_images_from_file(f, tmp_dir)
                self._log(summary)
                if imgs:
                    self._enqueue_files(imgs)

        if not image_files and not doc_files:
            self._log("⚠  Dropped items contained no recognized files")

    # ── Import file (PDF/DOCX/etc.) ──────────────────────────────────────────

    def _import_file(self):
        files = filedialog.askopenfilenames(title="Import File to Extract Images",
                                            filetypes=IMPORT_FORMATS)
        if not files:
            return
        tmp_dir = tempfile.mkdtemp(prefix="jpeg_imp_")
        self._tmp_dirs.append(tmp_dir)
        total = 0
        for filepath in files:
            self._log(f"🔍  Extracting: {Path(filepath).name} …")
            self.root.update_idletasks()
            imgs, summary = extract_images_from_file(filepath, tmp_dir)
            self._log(summary)
            if imgs:
                self._enqueue_files(imgs)
                total += len(imgs)
        if total:
            messagebox.showinfo("Import Complete",
                                f"{total} image(s) extracted and added.")
        else:
            messagebox.showwarning("No Images Found",
                                   "No extractable images were found.")

    # ── Output folder ────────────────────────────────────────────────────────

    def _browse_output(self):
        folder = filedialog.askdirectory(title="Select Output Folder")
        if folder:
            self.out_var.set(folder)
            self.out_entry.config(fg=THEMES[self._theme]["entry_fg"])

    def _clear_output(self):
        self.out_var.set("")
        self.out_entry.config(fg=THEMES[self._theme]["hint"])
        self.out_entry.delete(0, "end")
        self.out_entry.insert(0, "No folder selected — click Browse to choose…")

    def _open_output_folder(self):
        folder = self._last_output_folder or self.out_var.get()
        if folder and os.path.isdir(folder):
            open_folder(folder)
        else:
            messagebox.showinfo("No Folder", "No valid output folder to open yet.")

    # ── Settings ─────────────────────────────────────────────────────────────

    def _update_q_label(self, *_):
        q = self.quality_var.get()
        label = "Low" if q < 50 else "Medium" if q < 75 else "High" if q < 95 else "Maximum"
        self.q_label.config(text=f"{q} ({label})")

    def _update_prefix_hint(self, *_):
        p = self.prefix_var.get() or "Image"
        self.prefix_hint.config(text=f"→  files saved as {p}_1.jpg, {p}_2.jpg, …")

    # ── Logging ──────────────────────────────────────────────────────────────

    def _log(self, msg):
        self.log_box.insert("end", msg)
        self.log_box.yview("end")

    def _log_clear(self):
        self.log_box.delete(0, "end")

    def _set_progress(self, current, total):
        self.progress["maximum"] = total
        self.progress["value"] = current

    # ── Conversion (threaded) ────────────────────────────────────────────────

    def _convert(self):
        files = self._get_ordered_files()
        if not files:
            messagebox.showwarning("No Images", "Please add at least one image.")
            return

        output = self.out_var.get()
        if not output or not os.path.isdir(output):
            messagebox.showwarning("No Output Folder Selected",
                                   "Please select an output folder before converting.")
            return

        prefix = self.prefix_var.get().strip() or "Image"

        quality = self.quality_var.get()
        self._last_output_folder = output

        # UI state for conversion in progress
        self._converting = True
        self._cancel_event.clear()
        self.convert_btn.config(state="disabled")
        self.cancel_btn.config(state="normal")
        self.open_folder_btn.config(state="disabled")
        self._log_clear()
        self.progress["maximum"] = len(files)
        self.progress["value"] = 0

        threading.Thread(
            target=self._convert_worker,
            args=(files, output, prefix, quality),
            daemon=True,
        ).start()

    def _convert_worker(self, files, output, prefix, quality):
        used_names = set()
        ok = fail = 0
        cancelled = False

        for i, filepath in enumerate(files):
            if self._cancel_event.is_set():
                cancelled = True
                self.root.after(0, self._log, "⚠  Cancelled by user")
                break
            success, msg = convert_image_to_jpeg(
                filepath, output, quality, i + 1, prefix, used_names
            )
            self.root.after(0, self._log, msg)
            if success:
                ok += 1
            else:
                fail += 1
            self.root.after(0, self._set_progress, i + 1, len(files))

        self.root.after(0, self._convert_done, ok, fail, cancelled)

    def _convert_done(self, ok, fail, cancelled):
        self._converting = False
        self.convert_btn.config(state="normal")
        self.cancel_btn.config(state="disabled")
        self.open_folder_btn.config(state="normal")

        summary = f"Done — {ok} converted"
        if fail:
            summary += f", {fail} failed"
        if cancelled:
            summary += " (cancelled)"
        self._log("─" * 60)
        self._log(summary)

        if cancelled:
            messagebox.showinfo("Cancelled", summary)
        else:
            messagebox.showinfo("Conversion Complete", summary)

    def _cancel_convert(self):
        if self._converting:
            self._cancel_event.set()
            self._log("🛑 Cancel requested …")

    # ── Cleanup ──────────────────────────────────────────────────────────────

    def _on_close(self):
        for d in self._tmp_dirs:
            try:
                shutil.rmtree(d, ignore_errors=True)
            except Exception:
                pass
        self.root.destroy()


def main():
    root = TkinterDnD.Tk()
    ConverterApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()
